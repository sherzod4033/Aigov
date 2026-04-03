import asyncio
from playwright.async_api import async_playwright
import pptx
from pptx.util import Inches
import os

async def main():
    print("Starting Playwright to capture slides...")
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()
        
        # Set viewport to 16:9 1080p
        await page.set_viewport_size({"width": 1920, "height": 1080})
        
        file_path = f"file://{os.path.abspath('SoliqAI_Presentation.html')}"
        await page.goto(file_path)
        
        # Wait for reveal.js to fully load
        await page.wait_for_selector('.reveal', state='attached')
        await asyncio.sleep(3)  # wait for animations and fonts to settle
        
        # Disable controls, progress, and slide number so they don't appear in screenshots
        await page.evaluate('''
            document.querySelectorAll('.controls, .progress, .slide-number').forEach(e => e.style.display = 'none');
            // ensure any slide transitions are instant for capturing
            Reveal.configure({ transition: 'none' });
        ''')
        
        total_slides = await page.evaluate("Reveal.getTotalSlides()")
        print(f"Total slides found: {total_slides}")
        
        # Generate the pptx
        prs = pptx.Presentation()
        
        # Define 16:9 layout
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6]
        
        for i in range(total_slides):
            print(f"Capturing slide {i+1}...")
            
            # Go to the specific slide using Reveal API to ensure accuracy
            await page.evaluate(f"Reveal.slide({i})")
            await asyncio.sleep(1) # wait a moment for the slide to render correctly
            
            screenshot_path = f"slide_{i}.png"
            await page.screenshot(path=screenshot_path)
            
            # Add to PPTX
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(screenshot_path, 0, 0, width=Inches(16), height=Inches(9))
            
        prs.save("SoliqAI_Presentation.pptx")
        print("Presentation saved to SoliqAI_Presentation.pptx")
        
        # Cleanup screenshots
        for i in range(total_slides):
            try:
                os.remove(f"slide_{i}.png")
            except Exception as e:
                pass
                
        await browser.close()

if __name__ == "__main__":
    asyncio.run(main())
